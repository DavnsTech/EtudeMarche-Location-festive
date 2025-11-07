"""
PowerPoint presentation generator for market study findings
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import pandas as pd
import os
from typing import Optional, List, Dict, Any

class MarketStudyPresentation:
    def __init__(self, 
                 competitor_data_file: str = 'data/competitor_research.xlsx',
                 market_data_file: str = 'data/market_overview.xlsx'):
        """
        Initializes the PowerPoint presentation generator.

        Args:
            competitor_data_file (str): Path to the competitor research Excel file.
            market_data_file (str): Path to the market overview Excel file.
        """
        self.prs = Presentation()
        # Standard widescreen aspect ratio (16:9)
        self.prs.slide_width = Inches(13.33) 
        self.prs.slide_height = Inches(7.5)
        self.competitor_data_file = competitor_data_file
        self.market_data_file = market_data_file
        self.reports_dir: str = 'reports'
        self.output_filename: str = os.path.join(self.reports_dir, 'Location_Festive_Niort_Market_Study_Presentation.pptx')

        # Define color scheme
        self.primary_color = RGBColor(0x4F, 0x81, 0xBD) # Blue
        self.secondary_color = RGBColor(0xC0, 0x50, 0x4E) # Red
        self.accent_color = RGBColor(0x9B, 0x9B, 0x9B) # Gray

        # Ensure reports directory exists
        if not os.path.exists(self.reports_dir):
            try:
                os.makedirs(self.reports_dir)
                print(f"Created directory: {self.reports_dir}")
            except OSError as e:
                print(f"Error creating directory {self.reports_dir}: {e}")

    def _add_title_slide(self, title_text: str, subtitle_text: str) -> None:
        """Adds a title slide."""
        slide_layout = self.prs.slide_layouts[0]  # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text
        
        # Style the title
        title_format = title.text_frame.paragraphs[0]
        title_format.font.size = Pt(40)
        title_format.font.bold = True
        title_format.font.color.rgb = self.primary_color
        title_format.alignment = PP_ALIGN.CENTER
        
        subtitle_format = subtitle.text_frame.paragraphs[0]
        subtitle_format.font.size = Pt(24)
        subtitle_format.font.color.rgb = self.accent_color
        subtitle_format.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, title_text: str, content_items: List[str]) -> None:
        """Adds a title and content slide with bullet points."""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        
        title.text = title_text
        tf = body.text_frame
        tf.clear() # Clear existing content

        # Add bullet points
        for i, item in enumerate(content_items):
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x00, 0x00, 0x00) # Black text
            p.level = 0
            if i == 0: # First item for title styling
                p.font.bold = True
                p.font.color.rgb = self.primary_color

    def _load_excel_data(self, filepath: str) -> Optional[pd.DataFrame]:
        """Loads data from an Excel file, handling potential errors."""
        if not os.path.exists(filepath):
            print(f"Error: Data file not found at {filepath}")
            return None
        try:
            df = pd.read_excel(filepath)
            return df
        except Exception as e:
            print(f"Error reading Excel file {filepath}: {e}")
            return None

    def _add_chart_slide(self, title_text: str, chart_data_dict: Dict[str, Any], chart_type: str = 'bar') -> None:
        """Adds a slide with a chart."""
        slide_layout = self.prs.slide_layouts[5] # Blank slide layout for more control
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title shape
        left = top = width = height = Inches(0)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.primary_color
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Position chart
        chart_left = Inches(1)
        chart_top = Inches(1.5)
        chart_width = Inches(10)
        chart_height = Inches(5)

        if chart_type == 'bar':
            chart_data = CategoryChartData()
            chart_data.categories = chart_data_dict['categories']
            for series_name, values in chart_data_dict['series'].items():
                chart_data.add_series(series_name, values)

            x, y, cx, cy = chart_left, chart_top, chart_width, chart_height
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            chart = graphic_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.TOP_RIGHT
            chart.value_axis.has_major_gridlines = True
            chart.category_axis.has_major_gridlines = False # Often cleaner without vertical gridlines
            chart.category_axis.tick_labels.font.size = Pt(10)
            chart.value_axis.tick_labels.font.size = Pt(10)
            chart.plots[0].has_data_labels = True
            chart.plots[0].data_labels.number_format = "0" # Ensure integer display
            chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        elif chart_type == 'pie':
            chart_data = CategoryChartData()
            chart_data.categories = chart_data_dict['categories']
            for series_name, values in chart_data_dict['series'].items():
                chart_data.add_series(series_name, values)

            x, y, cx, cy = chart_left, chart_top, chart_width, chart_height
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            )
            chart = graphic_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.plots[0].has_data_labels = True
            chart.plots[0].data_labels.show_percentage = True
            chart.plots[0].data_labels.show_val = False
            chart.plots[0].data_labels.font.size = Pt(10)
            
    def generate_presentation(self) -> Optional[str]:
        """Generates the complete PowerPoint presentation."""
        print("Generating PowerPoint presentation...")

        # Slide 1: Title Slide
        self._add_title_slide("Étude de Marché - Location Festive Niort", "Analyse Complète du Secteur")

        # Slide 2: Introduction/Executive Summary
        market_df = self._load_excel_data(self.market_data_file)
        market_summary_items = []
        if market_df is not None:
            market_summary_items.append(f"Industrie: {market_df.loc[market_df['Category'] == 'Industry', 'Details'].iloc[0]}")
            market_summary_items.append(f"Localisation Principale: {market_df.loc[market_df['Category'] == 'Primary Location', 'Details'].iloc[0]}")
            market_summary_items.append(f"Tendances Clés: {market_df.loc[market_df['Category'] == 'Key Trends', 'Details'].iloc[0]}")
        else:
            market_summary_items.append("Données de marché indisponibles.")
        self._add_content_slide("Résumé Exécutif", market_summary_items)

        # Slide 3: Competitor Overview
        competitor_df = self._load_excel_data(self.competitor_data_file)
        competitor_summary_items = []
        if competitor_df is not None:
            total_competitors = len(competitor_df)
            avg_strengths = competitor_df['Strengths'].dropna().apply(lambda x: len(str(x).split(',')) if pd.notnull(x) and str(x).strip() else 0).mean()
            avg_weaknesses = competitor_df['Weaknesses'].dropna().apply(lambda x: len(str(x).split(',')) if pd.notnull(x) and str(x).strip() else 0).mean()
            
            competitor_summary_items.append(f"Nombre total de concurrents analysés : {total_competitors}")
            competitor_summary_items.append(f"Nombre moyen de points forts par concurrent : {avg_strengths:.2f}")
            competitor_summary_items.append(f"Nombre moyen de points faibles par concurrent : {avg_weaknesses:.2f}")
            
            market_pos_counts = competitor_df['Market Position'].value_counts()
            if not market_pos_counts.empty:
                 competitor_summary_items.append("Répartition du positionnement sur le marché :")
                 for position, count in market_pos_counts.items():
                     competitor_summary_items.append(f"  - {position}: {count}")
            else:
                competitor_summary_items.append("Positionnement sur le marché non spécifié pour la plupart des concurrents.")
        else:
            competitor_summary_items.append("Données concurrentielles indisponibles.")
        self._add_content_slide("Aperçu des Concurrents", competitor_summary_items)

        # Slide 4: Competitor Strengths vs. Weaknesses Bar Chart
        if competitor_df is not None and not competitor_df.empty:
            competitor_df['Num_Strengths'] = competitor_df['Strengths'].dropna().apply(
                lambda x: len(str(x).split(',')) if pd.notnull(x) and str(x).strip() else 0
            )
            competitor_df['Num_Weaknesses'] = competitor_df['Weaknesses'].dropna().apply(
                lambda x: len(str(x).split(',')) if pd.notnull(x) and str(x).strip() else 0
            )
            
            chart_data_dict = {
                'categories': competitor_df['Competitor'].tolist(),
                'series': {
                    'Forces': competitor_df['Num_Strengths'].tolist(),
                    'Faiblesses': competitor_df['Num_Weaknesses'].tolist()
                }
            }
            self._add_chart_slide("Comparaison des Forces et Faiblesses par Concurrent", chart_data_dict, chart_type='bar')
        else:
            self._add_content_slide("Analyse des Forces et Faiblesses", ["Données concurrentielles indisponibles pour générer le graphique."])

        # Slide 5: Market Position Pie Chart
        if competitor_df is not None and not competitor_df.empty:
            market_pos_counts = competitor_df['Market Position'].value_counts()
            if not market_pos_counts.empty:
                chart_data_dict = {
                    'categories': market_pos_counts.index.tolist(),
                    'series': {
                        'Nombre de Concurrents': market_pos_counts.values.tolist()
                    }
                }
                self._add_chart_slide("Répartition du Positionnement sur le Marché", chart_data_dict, chart_type='pie')
            else:
                self._add_content_slide("Positionnement sur le Marché", ["Aucune donnée de positionnement sur le marché disponible."])
        else:
             self._add_content_slide("Positionnement sur le Marché", ["Données concurrentielles indisponibles pour générer le graphique."])

        # Slide 6: Market Trends
        market_trends_items = []
        if market_df is not None:
            trends_data = market_df[market_df['Category'] == 'Key Trends']
            if not trends_data.empty:
                trends_str = trends_data['Details'].iloc[0]
                if pd.notna(trends_str):
                    market_trends_items = [f"- {item.strip()}" for item in trends_str.split(',')]
            if not market_trends_items:
                market_trends_items.append("Aucune tendance de marché spécifique identifiée.")
        else:
            market_trends_items.append("Données de marché indisponibles.")
        self._add_content_slide("Tendances Clés du Marché", market_trends_items)

        # Slide 7: Seasonality
        seasonality_items = []
        if market_df is not None:
            seasonality_data = market_df[market_df['Category'] == 'Seasonal Peaks']
            if not seasonality_data.empty:
                seasonality_str = seasonality_data['Details'].iloc[0]
                if pd.notna(seasonality_str):
                    seasonality_items = [f"- {item.strip()}" for item in seasonality_str.split(',')]
            if not seasonality_items:
                seasonality_items.append("Aucune information sur la saisonnalité du marché identifiée.")
        else:
            seasonality_items.append("Données de marché indisponibles.")
        self._add_content_slide("Facteurs de Saisonnalité", seasonality_items)
        
        # Slide 8: Conclusion/Recommendations (Placeholder)
        self._add_content_slide("Conclusion et Recommandations", [
            "Synthèse des analyses menées.",
            "Identification des opportunités stratégiques.",
            "Recommandations pour 'Location Festive Niort'.",
            "(Cette section nécessite une analyse plus approfondie des données.)"
        ])

        # Save the presentation
        try:
            self.prs.save(self.output_filename)
            print(f"PowerPoint presentation saved successfully to: {self.output_filename}")
            return self.output_filename
        except Exception as e:
            print(f"Error saving PowerPoint presentation: {e}")
            return None

